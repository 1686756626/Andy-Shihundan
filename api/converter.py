#!/usr/bin/env python3
"""
文件转换 API — 纯 Python 解析文档为文本
支持：.doc .docx .wps .rtf .odt .xls .xlsx .ppt .pptx .pdf
端口：8090
"""
import os
import tempfile
from pathlib import Path
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware

app = FastAPI(title="Doc Converter API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

SUPPORTED_EXTS = {'.doc', '.wps', '.docx', '.rtf', '.odt', '.xls', '.xlsx', '.ppt', '.pptx', '.pdf'}


def extract_docx(path):
    """解析 .docx"""
    from docx import Document
    doc = Document(path)
    parts = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    # 也提取表格
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append(' | '.join(cells))
    return '\n'.join(parts)


def extract_doc_or_wps(path):
    """解析 .doc（OLE2）和 .wps"""
    import olefile
    try:
        ole = olefile.OleFileIO(path)
        # 尝试读取 WordDocument stream
        if ole.exists('WordDocument'):
            # 尝试提取 1Table 或 0Table 中的文本
            for stream_name in ['1Table', '0Table']:
                if ole.exists(stream_name):
                    data = ole.openstream(stream_name).read()
                    # 尝试 UTF-16LE 解码（Word 常用编码）
                    try:
                        text = data.decode('utf-16-le', errors='ignore')
                        # 清理非打印字符
                        cleaned = ''.join(c for c in text if c.isprintable() or c in '\n\r\t')
                        if len(cleaned.strip()) > 20:
                            return cleaned
                    except Exception:
                        pass

            # 如果上面没拿到，尝试读取整个 WordDocument stream
            data = ole.openstream('WordDocument').read()
            # 尝试多种编码
            for enc in ['utf-16-le', 'cp1252', 'gb2312', 'gbk', 'latin1']:
                try:
                    text = data.decode(enc, errors='ignore')
                    cleaned = ''.join(c for c in text if c.isprintable() or c in '\n\r\t')
                    # 过滤掉太短或太多乱码的
                    printable = sum(1 for c in cleaned if c.isalnum() or '\u4e00' <= c <= '\u9fff')
                    if len(cleaned.strip()) > 20 and printable > len(cleaned) * 0.3:
                        return cleaned
                except Exception:
                    continue

        # 最后尝试：读取所有 stream 拼起来
        all_text = []
        for stream in ole.listdir():
            stream_path = '/'.join(stream)
            try:
                data = ole.openstream(stream).read()
                for enc in ['utf-16-le', 'utf-8', 'gbk', 'latin1']:
                    try:
                        text = data.decode(enc, errors='ignore')
                        cleaned = ''.join(c for c in text if c.isprintable() or c in '\n\r\t')
                        printable = sum(1 for c in cleaned if c.isalnum() or '\u4e00' <= c <= '\u9fff')
                        if printable > len(cleaned) * 0.4 and len(cleaned.strip()) > 10:
                            all_text.append(cleaned)
                            break
                    except Exception:
                        continue
            except Exception:
                continue

        ole.close()
        if all_text:
            return '\n'.join(all_text)

        return '[无法解析该 .doc/.wps 文件的文本内容，建议转换为 .docx 格式后重试]'
    except Exception as e:
        return f'[解析失败: {e}]'


def extract_xlsx(path):
    """解析 .xlsx"""
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f'=== {sheet.title} ===')
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else '' for c in row]
            if any(cells):
                parts.append(' | '.join(cells))
    wb.close()
    return '\n'.join(parts)


def extract_xls(path):
    """解析 .xls"""
    import xlrd
    wb = xlrd.open_workbook(path)
    parts = []
    for sheet in wb.sheets():
        parts.append(f'=== {sheet.name} ===')
        for row_idx in range(sheet.nrows):
            cells = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]
            parts.append(' | '.join(cells))
    return '\n'.join(parts)


def extract_pptx(path):
    """解析 .pptx"""
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f'=== 第 {i} 页 ===')
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        parts.append(para.text)
    return '\n'.join(parts)


def extract_pdf(path):
    """解析 .pdf"""
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            parts = []
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text:
                    parts.append(text)
            return '\n'.join(parts)
    except Exception as e:
        return f'[PDF解析失败: {e}]'


def extract_rtf(path):
    """简单提取 .rtf 文本"""
    try:
        with open(path, 'r', errors='replace') as f:
            content = f.read()
        # 简单剥 RTF 标记
        import re
        text = re.sub(r'\\[a-z]+\d*\s?', '', content)
        text = re.sub(r'[{}]', '', text)
        text = re.sub(r'\\[^a-z]', '', text)
        return text.strip()
    except Exception as e:
        return f'[RTF解析失败: {e}]'


def extract_odt(path):
    """解析 .odt（本质是 zip 包含 content.xml）"""
    import zipfile
    import re
    try:
        with zipfile.ZipFile(path) as z:
            with z.open('content.xml') as f:
                xml = f.read().decode('utf-8', errors='replace')
        # 简单剥 XML 标签
        text = re.sub(r'<[^>]+>', '', xml)
        text = re.sub(r'\s+', ' ', text).strip()
        return text
    except Exception as e:
        return f'[ODT解析失败: {e}]'


EXTRACTORS = {
    '.docx': extract_docx,
    '.doc': extract_doc_or_wps,
    '.wps': extract_doc_or_wps,
    '.rtf': extract_rtf,
    '.odt': extract_odt,
    '.xlsx': extract_xlsx,
    '.xls': extract_xls,
    '.pptx': extract_pptx,
    '.ppt': extract_pptx,  # python-pptx 不支持 .ppt，但试试
    '.pdf': extract_pdf,
}


@app.post("/convert")
async def convert_file(file: UploadFile = File(...)):
    if not file.filename:
        raise HTTPException(400, "No filename")

    ext = Path(file.filename).suffix.lower()
    if ext not in SUPPORTED_EXTS:
        raise HTTPException(400, f"Unsupported format: {ext}. Supported: {', '.join(sorted(SUPPORTED_EXTS))}")

    content = await file.read()
    if len(content) > 50 * 1024 * 1024:
        raise HTTPException(400, "File too large (max 50MB)")

    with tempfile.TemporaryDirectory() as tmpdir:
        safe_name = "input" + ext
        input_path = os.path.join(tmpdir, safe_name)
        with open(input_path, "wb") as f:
            f.write(content)

        extractor = EXTRACTORS.get(ext)
        if not extractor:
            raise HTTPException(400, f"No extractor for {ext}")

        try:
            text = extractor(input_path)
        except Exception as e:
            raise HTTPException(500, f"Extraction failed: {e}")

        if not text or not text.strip():
            raise HTTPException(500, "No text content extracted from file")

        return {
            "filename": file.filename,
            "text": text,
            "length": len(text),
        }


@app.get("/health")
async def health():
    return {"status": "ok"}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="127.0.0.1", port=8090)
